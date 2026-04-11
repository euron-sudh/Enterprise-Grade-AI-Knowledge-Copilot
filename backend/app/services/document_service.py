"""
Document service — handles file upload, text extraction, and chunking.
"""
import hashlib
import logging
import os
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional, Tuple

from sqlalchemy import delete, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import settings
from app.models.knowledge import Document, DocumentChunk, DocumentStatus
from app.models.user import User
from app.services.vector_store import get_vector_store

logger = logging.getLogger(__name__)

CONNECTOR_LIKE_TYPES = {
    "google_drive",
    "gmail",
    "github",
    "gitlab",
    "confluence",
    "notion",
    "slack",
    "jira",
    "salesforce",
    "hubspot",
    "zendesk",
    "intercom",
    "sharepoint",
    "onedrive",
    "dropbox",
    "web",
    "web_crawler",
    "figma",
}

SUPPORTED_TYPES = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".txt": "txt",
    ".md": "md",
    ".csv": "csv",
    ".json": "json",
    ".html": "html",
    ".htm": "html",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".webp": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
}

IMAGE_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
}

CHUNK_SIZE_TOKENS = 500
CHUNK_OVERLAP_TOKENS = 50


def _estimate_tokens(text: str) -> int:
    """Rough token estimate: ~4 chars per token."""
    return max(1, len(text) // 4)


def build_content_duplicate_key(file_type: str, file_content: bytes) -> str:
    digest = hashlib.sha256(file_content).hexdigest()
    return f"content:{file_type}:{digest}"


def build_path_duplicate_key(file_type: str, file_path: str) -> str:
    normalized_path = (file_path or "").strip().lower()
    return f"path:{file_type}:{normalized_path}"


def build_duplicate_key(
    file_type: str,
    *,
    file_content: bytes | None = None,
    file_path: str | None = None,
    original_name: str | None = None,
    file_size: int | None = None,
) -> str:
    if file_content is not None and file_type not in CONNECTOR_LIKE_TYPES:
        return build_content_duplicate_key(file_type, file_content)
    if file_path:
        return build_path_duplicate_key(file_type, file_path)

    fallback_basis = "|".join(
        [file_type, original_name or "", str(file_size or 0)]
    )
    return f"fallback:{hashlib.sha256(fallback_basis.encode('utf-8')).hexdigest()}"


async def find_duplicate_document(
    db: AsyncSession,
    *,
    user_id: uuid.UUID,
    duplicate_key: str,
) -> Document | None:
    result = await db.execute(
        select(Document)
        .where(Document.user_id == user_id, Document.duplicate_key == duplicate_key)
        .limit(1)
    )
    return result.scalar_one_or_none()


def derive_existing_duplicate_key(document: Document) -> str:
    if document.duplicate_key:
        return document.duplicate_key

    if document.file_type in CONNECTOR_LIKE_TYPES or "://" in (document.file_path or ""):
        return build_duplicate_key(
            document.file_type,
            file_path=document.file_path,
            original_name=document.original_name,
            file_size=document.file_size,
        )

    path = Path(document.file_path) if document.file_path else None
    if path and path.is_file():
        return build_duplicate_key(
            document.file_type,
            file_content=path.read_bytes(),
            original_name=document.original_name,
            file_size=document.file_size,
        )

    return build_duplicate_key(
        document.file_type,
        original_name=document.original_name,
        file_size=document.file_size,
    )


def _document_rank(document: Document) -> tuple[int, int, int, datetime]:
    status_value = document.status.value if hasattr(document.status, "value") else str(document.status)
    status_rank = {
        DocumentStatus.indexed.value: 3,
        DocumentStatus.processing.value: 2,
        DocumentStatus.failed.value: 1,
    }.get(status_value, 0)
    return (
        status_rank,
        document.word_count or 0,
        document.file_size or 0,
        document.updated_at or document.created_at,
    )


async def cleanup_duplicate_documents(db: AsyncSession, user_id: uuid.UUID | None = None) -> dict[str, int]:
    query = select(Document).order_by(Document.created_at.desc())
    if user_id is not None:
        query = query.where(Document.user_id == user_id)

    result = await db.execute(query)
    documents = list(result.scalars().all())

    kept_by_key: dict[tuple[uuid.UUID, str], Document] = {}
    duplicate_ids: list[uuid.UUID] = []
    duplicate_files: list[str] = []
    updated_keys = 0

    for document in documents:
        duplicate_key = derive_existing_duplicate_key(document)
        if document.duplicate_key != duplicate_key:
            document.duplicate_key = duplicate_key
            updated_keys += 1

        group_key = (document.user_id, duplicate_key)
        current = kept_by_key.get(group_key)
        if current is None:
            kept_by_key[group_key] = document
            continue

        keep_document = current if _document_rank(current) >= _document_rank(document) else document
        duplicate_document = document if keep_document is current else current
        kept_by_key[group_key] = keep_document
        duplicate_ids.append(duplicate_document.id)
        if duplicate_document.file_path and "://" not in duplicate_document.file_path:
            duplicate_files.append(duplicate_document.file_path)

    if duplicate_ids:
        await db.execute(delete(DocumentChunk).where(DocumentChunk.document_id.in_(duplicate_ids)))
        await db.execute(delete(Document).where(Document.id.in_(duplicate_ids)))

    await db.flush()

    for file_path in duplicate_files:
        try:
            path = Path(file_path)
            if path.is_file():
                path.unlink(missing_ok=True)
        except Exception as exc:
            logger.warning("Failed to remove duplicate file %s: %s", file_path, exc)

    return {
        "scanned": len(documents),
        "updated_keys": updated_keys,
        "deleted_duplicates": len(duplicate_ids),
    }


def _chunk_text(text: str) -> List[Tuple[int, str]]:
    """
    Split text into overlapping chunks.
    Returns list of (chunk_index, chunk_text).
    """
    # Convert token targets to rough char counts
    chunk_size_chars = CHUNK_SIZE_TOKENS * 4
    overlap_chars = CHUNK_OVERLAP_TOKENS * 4

    chunks = []
    start = 0
    idx = 0

    while start < len(text):
        end = start + chunk_size_chars

        # Try to break at a sentence boundary
        if end < len(text):
            # Look for sentence end within last 20% of chunk
            boundary_search_start = start + int(chunk_size_chars * 0.8)
            for sep in [". ", ".\n", "! ", "? ", "\n\n"]:
                pos = text.rfind(sep, boundary_search_start, end)
                if pos != -1:
                    end = pos + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append((idx, chunk))
            idx += 1

        # Move forward with overlap
        start = end - overlap_chars
        if start >= len(text):
            break

    return chunks


def _extract_text_pdf(file_path: str) -> Tuple[str, int]:
    """Extract text from PDF using PyPDF2. Returns (text, page_count)."""
    try:
        import PyPDF2

        text_parts = []
        page_count = 0
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            page_count = len(reader.pages)
            for page in reader.pages:
                try:
                    text_parts.append(page.extract_text() or "")
                except Exception:
                    pass
        return "\n".join(text_parts), page_count
    except Exception as e:
        logger.warning(f"PDF extraction failed: {e}")
        return "", 0


def _extract_text_docx(file_path: str) -> Tuple[str, int]:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document as DocxDoc

        doc = DocxDoc(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        # Estimate page count (~500 words per page)
        word_count = len(text.split())
        page_count = max(1, word_count // 500)
        return text, page_count
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return "", 0


def _extract_text_plain(file_path: str) -> Tuple[str, int]:
    """Read plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        word_count = len(text.split())
        page_count = max(1, word_count // 500)
        return text, page_count
    except Exception as e:
        logger.warning(f"Plain text extraction failed: {e}")
        return "", 0


def _sanitize_text(text: str) -> str:
    """Remove characters PostgreSQL UTF8 can't store (null bytes, etc.)."""
    # Remove null bytes — PostgreSQL VARCHAR rejects \x00
    return text.replace('\x00', '').replace('\x0b', ' ').replace('\x0c', ' ')


def _extract_text_xlsx(file_path: str) -> Tuple[str, int]:
    """Extract text from Excel workbook using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append("\t".join(cells))
        wb.close()
        text = "\n".join(parts)
        return text, max(1, len(wb.sheetnames))
    except Exception as e:
        logger.warning(f"XLSX extraction failed: {e}")
        return "", 0


def _extract_text_pptx(file_path: str) -> Tuple[str, int]:
    """Extract text from PowerPoint using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        text = "\n".join(parts)
        return text, len(prs.slides)
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return "", 0


async def _extract_text_image(file_path: str, file_ext: str) -> Tuple[str, int]:
    """
    Use Claude vision API to generate a description of an image.
    Falls back to a placeholder if the API key is not configured.
    """
    import base64

    try:
        from app.config import settings
        api_key = getattr(settings, "ANTHROPIC_API_KEY", "") or ""
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY not configured")

        import anthropic
        with open(file_path, "rb") as f:
            image_bytes = f.read()

        media_type = IMAGE_MIME_TYPES.get(file_ext, "image/png")
        # Supported media types for Anthropic vision
        if media_type not in ("image/png", "image/jpeg", "image/gif", "image/webp"):
            media_type = "image/png"

        b64_data = base64.standard_b64encode(image_bytes).decode("utf-8")
        client = anthropic.Anthropic(api_key=api_key)
        message = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1024,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": b64_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": (
                                "Please describe the content of this image in detail. "
                                "Include any text visible in the image, the main subject, "
                                "charts or diagrams if present, and any other relevant information. "
                                "Write your response as plain text suitable for a knowledge base."
                            ),
                        },
                    ],
                }
            ],
        )
        description = message.content[0].text if message.content else ""
        if not description:
            raise ValueError("Empty response from vision API")

        file_name = Path(file_path).name
        text = f"[Image: {file_name}]\n\n{description}"
        return _sanitize_text(text), 1

    except Exception as e:
        logger.warning(f"Image vision extraction failed for {file_path}: {e}")
        file_name = Path(file_path).name
        placeholder = (
            f"[Image file: {file_name}]\n\n"
            f"This is an image file that could not be automatically analyzed. "
            f"File size: {Path(file_path).stat().st_size} bytes."
        )
        return placeholder, 1


def extract_text(file_path: str, file_ext: str) -> Tuple[str, int]:
    """Dispatch extraction based on file extension. For images, use async variant."""
    if file_ext == ".pdf":
        text, pages = _extract_text_pdf(file_path)
    elif file_ext in (".docx", ".doc"):
        text, pages = _extract_text_docx(file_path)
    elif file_ext in (".xlsx", ".xls"):
        text, pages = _extract_text_xlsx(file_path)
    elif file_ext == ".pptx":
        text, pages = _extract_text_pptx(file_path)
    elif file_ext in IMAGE_MIME_TYPES:
        # Image files — return empty so process_document uses the async path
        return "", 1
    else:
        text, pages = _extract_text_plain(file_path)
    return _sanitize_text(text), pages


async def save_uploaded_file(
    file_content: bytes,
    original_filename: str,
    user_id: str,
) -> Tuple[str, str, str]:
    """
    Save raw file bytes to disk.
    Returns (saved_path, file_ext, file_type_label).
    """
    upload_dir = Path(settings.UPLOAD_DIR) / str(user_id)
    upload_dir.mkdir(parents=True, exist_ok=True)

    ext = Path(original_filename).suffix.lower()
    file_type = SUPPORTED_TYPES.get(ext, "unknown")

    unique_name = f"{uuid.uuid4()}{ext}"
    saved_path = upload_dir / unique_name

    with open(saved_path, "wb") as f:
        f.write(file_content)

    return str(saved_path), ext, file_type


async def process_document(
    document: Document,
    file_path: str,
    file_ext: str,
    db: AsyncSession,
) -> None:
    """
    Extract text from a document, create chunks, and mark it as indexed.
    """
    try:
        if file_ext in IMAGE_MIME_TYPES:
            text, page_count = await _extract_text_image(file_path, file_ext)
        else:
            text, page_count = extract_text(file_path, file_ext)

        word_count = len(text.split()) if text else 0
        chunks = _chunk_text(text) if text else []

        # Create chunk records
        chunk_records = []
        for chunk_idx, chunk_content in chunks:
            token_count = _estimate_tokens(chunk_content)
            chunk_record = DocumentChunk(
                document_id=document.id,
                content=chunk_content,
                chunk_index=chunk_idx,
                token_count=token_count,
            )
            db.add(chunk_record)
            chunk_records.append((chunk_idx, chunk_content, chunk_record))

        # Update document metadata
        document.status = DocumentStatus.indexed
        document.page_count = page_count
        document.word_count = word_count
        document.updated_at = datetime.now(timezone.utc)

        await db.flush()

        # Upsert to Pinecone vector store for semantic search
        vector_store = get_vector_store()
        if vector_store.enabled and chunk_records:
            try:
                vector_batch = []
                for chunk_idx, chunk_content, chunk_rec in chunk_records:
                    vector_batch.append({
                        "id": f"{document.id}:{chunk_idx}",
                        "text": chunk_content,
                        "metadata": {
                            "user_id": str(document.user_id),
                            "document_id": str(document.id),
                            "document_name": document.original_name or document.name,
                            "document_type": document.file_type or "upload",
                            "chunk_text": chunk_content[:2000],
                        },
                    })
                await vector_store.upsert_chunks(
                    user_id=document.user_id, records=vector_batch
                )
                logger.info(
                    f"Document {document.id} vectorized: {len(vector_batch)} chunks upserted to Pinecone"
                )
            except Exception as vec_err:
                logger.warning(
                    f"Pinecone upsert failed for document {document.id}: {vec_err}"
                )

        logger.info(
            f"Document {document.id} indexed: {len(chunks)} chunks, "
            f"{word_count} words, {page_count} pages"
        )

    except Exception as e:
        logger.error(f"Document processing failed for {document.id}: {e}")
        document.status = DocumentStatus.failed
        document.updated_at = datetime.now(timezone.utc)
        await db.flush()


async def upload_documents(
    files: List,  # List of UploadFile objects
    user: User,
    collection_id: Optional[str],
    tags: Optional[str],
    db: AsyncSession,
) -> List[Document]:
    """Handle multiple file uploads."""
    from fastapi import HTTPException, status

    tag_list = []
    if tags:
        tag_list = [t.strip() for t in tags.split(",") if t.strip()]

    documents = []
    max_size = settings.max_upload_size_bytes

    for upload_file in files:
        # Read file content
        content = await upload_file.read()

        if len(content) > max_size:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File '{upload_file.filename}' exceeds maximum size of {settings.MAX_UPLOAD_SIZE_MB}MB",
            )

        original_name = upload_file.filename or "unknown"
        file_ext = Path(original_name).suffix.lower()
        file_type = SUPPORTED_TYPES.get(file_ext, "unknown")
        duplicate_key = build_duplicate_key(
            file_type,
            file_content=content,
            original_name=original_name,
            file_size=len(content),
        )
        existing_document = await find_duplicate_document(
            db,
            user_id=user.id,
            duplicate_key=duplicate_key,
        )
        if existing_document is not None:
            # For image files that were previously stored with binary/garbage content,
            # re-process them so they get a proper AI-generated description.
            if file_ext in IMAGE_MIME_TYPES and existing_document.file_type == "image":
                result = await db.execute(
                    select(DocumentChunk)
                    .where(DocumentChunk.document_id == existing_document.id)
                    .limit(1)
                )
                first_chunk = result.scalar_one_or_none()
                is_binary_chunk = (
                    first_chunk is not None
                    and first_chunk.content
                    and not first_chunk.content.startswith("[Image")
                    and len([c for c in first_chunk.content[:100] if ord(c) > 127]) > 10
                )
                if is_binary_chunk:
                    logger.info(
                        "Re-processing image '%s' (document_id=%s) — previous content was binary",
                        original_name,
                        existing_document.id,
                    )
                    # Delete bad chunks
                    from sqlalchemy import delete as sa_delete
                    await db.execute(
                        sa_delete(DocumentChunk).where(DocumentChunk.document_id == existing_document.id)
                    )
                    await db.flush()
                    # Re-save the file and reprocess
                    file_path, file_ext_saved, file_type = await save_uploaded_file(
                        content, original_name, str(user.id)
                    )
                    existing_document.file_path = file_path
                    existing_document.status = DocumentStatus.processing
                    await db.flush()
                    await process_document(existing_document, file_path, file_ext, db)
            else:
                logger.info(
                    "Skipping duplicate upload '%s' for user %s (document_id=%s)",
                    original_name,
                    user.id,
                    existing_document.id,
                )
            documents.append(existing_document)
            continue

        file_path, file_ext, file_type = await save_uploaded_file(
            content, original_name, str(user.id)
        )

        # Create DB record
        document = Document(
            user_id=user.id,
            collection_id=uuid.UUID(collection_id) if collection_id else None,
            name=original_name,
            original_name=original_name,
            file_path=file_path,
            duplicate_key=duplicate_key,
            file_type=file_type,
            file_size=len(content),
            status=DocumentStatus.processing,
            tags=tag_list,
        )
        db.add(document)
        await db.flush()

        # Process synchronously (extract text and create chunks)
        await process_document(document, file_path, file_ext, db)

        documents.append(document)

    return documents
